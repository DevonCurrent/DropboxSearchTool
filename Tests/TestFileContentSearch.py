import unittest
import warnings
import sys, os

runPath = os.path.dirname(os.path.realpath(__file__))
sys.path.append(os.path.join(runPath, ".."))

import dropbox
from dropbox.exceptions import AuthError
from io import BytesIO
import tempfile

from docx import Document
from pptx import Presentation
import openpyxl
import PyPDF2
from subprocess import check_output

#from FileContentSearch import FileContentSearch


class TestFileContentSearch(unittest.TestCase):

    def test_docx(self):
        
        warnings.simplefilter("ignore", ResourceWarning)
        dropboxToken = open("DropboxSearchTokens.txt").readlines()[1].strip()
        
        dbx = dropbox.Dropbox(dropboxToken)

        try:
            dbx.users_get_current_account()
        except AuthError as err:
            print(err)
            print("ERROR: Invalid access token; try re-generating an access token from the app console on the web.")
            exit()

        try:
            metadata, f = dbx.files_download('/2014/Google/coolmoney.docx')
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        doc = Document(stream)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)

        #print('\n'.join(fullText))

    def test_doc(self):
        warnings.simplefilter("ignore", ResourceWarning)
        dropboxToken = open("DropboxSearchTokens.txt").readlines()[1].strip()
        
        dbx = dropbox.Dropbox(dropboxToken)

        out = ''

        with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as temp:
            dbx.files_download_to_file(temp.name, '/2014/Google/SRS4.doc')
            out = check_output(['antiword', '-f' ,temp.name])

        #print(out)


    def test_pptx(self):

        warnings.simplefilter("ignore", ResourceWarning)
        dropboxToken = open("DropboxSearchTokens.txt").readlines()[1].strip()
        
        dbx = dropbox.Dropbox(dropboxToken)

        try:
            dbx.users_get_current_account()
        except AuthError as err:
            print(err)
            print("ERROR: Invalid access token; try re-generating an access token from the app console on the web.")
            exit()

        try:
            metadata, f = dbx.files_download('/2014/Google/powertest.pptx')
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        #pptx only works with .pptx but not with .ppt
        prs = Presentation(stream)
        textRuns = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        textRuns.append(run.text)

        #print(textRuns)
    
    def test_xlsx(self):

        warnings.simplefilter("ignore", ResourceWarning)
        dropboxToken = open("DropboxSearchTokens.txt").readlines()[1].strip()
        
        dbx = dropbox.Dropbox(dropboxToken)

        try:
            dbx.users_get_current_account()
        except AuthError as err:
            print(err)
            print("ERROR: Invalid access token; try re-generating an access token from the app console on the web.")
            exit()

        try:
            metadata, f = dbx.files_download('/2014/Google/CapstoneExcel.xlsx')
        except dropbox.files.DownloadError as err:
            print(err)
            exit()
        
        stream = BytesIO(f.content)
        book = openpyxl.load_workbook(stream)

        string = ''

        for sheet in book:
            for row in sheet.iter_rows():
                for cell in row:
                    string += ' '
                    string += str(cell.value)

        #print(string)

    def test_pdf(self):
        warnings.simplefilter("ignore", ResourceWarning)
        dropboxToken = open("DropboxSearchTokens.txt").readlines()[1].strip()
        
        dbx = dropbox.Dropbox(dropboxToken)

        try:
            dbx.users_get_current_account()
        except AuthError as err:
            print(err)
            print("ERROR: Invalid access token; try re-generating an access token from the app console on the web.")
            exit()

        try:
            metadata, f = dbx.files_download('/2014/Google/databases_week_3_ER_modeling.pdf')
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        pdfReader = PyPDF2.PdfFileReader(stream)
        numberOfPages = pdfReader.numPages

        string = ''
        for i in range(0, numberOfPages):
            string += ' '
            string += str(pdfReader.getPage(i).extractText())

        #print(string)

if __name__ == "__main__":
    unittest.main()