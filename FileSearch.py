from BagOfWords import BagOfWords
import pdb

import dropbox
from dropbox.exceptions import AuthError

from io import BytesIO
import tempfile
from docx import Document
from pptx import Presentation
import openpyxl
import PyPDF2
from subprocess import check_output

import time

import concurrent.futures
"""
        Formats the fileList found on Dropbox to a list of each files' content. This is then passed to the 
        BagOfWords to find the most accurate searches.'

        Parameters
        ----------
        dropboxBot : class 'DropboxBot.DropboxBot'
            an instance of DropboxBot that has access to the Dropbox account
        fileList : list
            A list of files found on the Dropbox that are located in the specified companies and year fields.
        search : class 'Search.Search'
            Search object that contains tuples for keywords, companies, years, and specified searches by the Slack user
        type : 0 for content search, 1 for name search
        
        Returns
        -------
        accurateDocList
            The list of files that are most accurate to the search that the Slack user requested.
        """
class FileSearch:
    """
    Class used to represent the file search
    -----
    Methods
    -----
    file_search(self, dropboxBot, fileList, search)
        Formats the fileList found on Dropbox to a list of each files' content. This is then passed to the 
        BagOfWords to find the most accurate searches.
    _name_search(self, dropboxBot, fileList, search)
        Searches throught the file list based on file names
    _docx_parse(self, filename)
        Parses through Docx Files
    _doc_parse(self, filename)
        Parses through Doc Files
    _pptx_parse(self, filename)
        Parses through Pptx Files
    _xlsx_parse(self, filename)
        Parses through Xlsx Files
    _pdf_parse(self, filename)
        Parses through Pdf Files
    _nonsupported_parse(self, filename)
        Takes care of nonsupported files
    file_search(self, fileList, search, type)
        Searches through the fileList from Dropbox    
    """

    dropboxBot = None

    def _name_search(self, dropboxBot, fileList, search):
        """
    Searches throught the file list based on file names

    Parameters
    ----------
    dropboxBot : class 'DropboxBot.DropboxBot'
        an instance of DropboxBot that has access to the Dropbox account
    fileList : list
            A list of files found on the Dropbox that are located in the specified companies and year fields.
    search : class 'Search.Search'
            Search object that contains tuples for keywords, companies, years, and specified searches by the Slack user
    Returns
    -------
    BagOfWords.find_accurate_docs(fileList, fileNameTextList, keywords)
        The most accurate documents based on the BagOfWords class
    """
        fileNameTextList = []
        for entry in fileList:
            # this removes the .ext from entry names when adding to fileNameTextList
            splitEntry = entry.name.split('.')
            fileName = ""
            for word in splitEntry[:-1]: # there could be more than one '.'
                fileName += word
            fileNameTextList.append(fileName)

        keywords = ' '.join(search.keywords)

        return BagOfWords.find_accurate_docs(fileList, fileNameTextList, keywords)

    def _docx_parse(self, filename):
        """
    Parses through Docx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    '\n'.join(fullText)
        Full text of the parsed file
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        doc = Document(stream)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text.lower())

        return '\n'.join(fullText)

    def _doc_parse(self, filename):
        """
    Parses through Docx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    out.lower()
        the output of the parsed file in lowercase
    """
        out = ''
        with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as temp:
            self.dropboxBot.dbx.files_download_to_file(temp.name, filename)
            out = check_output(['antiword', '-f' ,temp.name])
        
        return out.lower()

    def _pptx_parse(self, filename):
        """
    Parses through Pptx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    text
        the output of the parsed file
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        #pptx only works with .pptx but not with .ppt
        prs = Presentation(stream)
        textRuns = []
            
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        textRuns.append(run.text.lower())

        text = ''
        for run in textRuns:
            text += run
        
        return text

    def _xlsx_parse(self, filename):
        """
    Parses through Xlsx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    string
        the output of the parsed file in lowercase
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        stream = BytesIO(f.content)
        book = openpyxl.load_workbook(stream)

        string = ''

        for sheet in book:
            for row in sheet.iter_rows():
                for cell in row:
                    string += ' '
                    string += str(cell.value).lower()

        return string

    def _pdf_parse(self, filename):
        """
    Parses through pdf Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    string
        the output of the parsed file in lowercase
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        pdfReader = PyPDF2.PdfFileReader(stream)

        numberOfPages = pdfReader.numPages
        string = ''

        for i in range(0, numberOfPages):
            string += ' '
            string += str(pdfReader.getPage(i).extractText().lower())

        return string

    def _nonsupported_parse(self, filename):
        """
    Takes care of nonsupported files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    fileName
        the file name
    """
        splitEntry = filename.split('.')
        fileName = ""

        for word in splitEntry[:-1]: # there could be more than one '.'
            fileName += word

        return fileName

    def Search(self, fileList, search, type):
        """
    Searches through the fileList from Dropbox  

    Parameters
    ----------
    fileList : list
            A list of files found on the Dropbox that are located in the specified companies and year fields.
    search : class 'Search.Search'
            Search object that contains tuples for keywords, companies, years, and specified searches by the Slack user
    type: int
            type of search
    Returns
    -------
    BagOfWords.find_accurate_docs(fileList, list, keywords)
        The most accurate documents based on the BagOfWords class
        
    """
        if type:
            return self._name_search(self.dropboxBot, fileList, search)

        list = []

        total1 = time.time()
        
        fileParseType = []
        """
        for file in fileList:
            filePath = file.path_display
            fileType = file.name.split('.')[-1]
            data = (fileType, filePath)
            fileParseType.append(data)
        """

        for index, file in enumerate(fileList, start=0):
            filePath = file.path_display
            fileType = file.name.split('.')[-1]
            data = (fileType, filePath, index)
            fileParseType.append(data)

        def download_and_parse(fileType, filePath, index):
            if (fileType == 'doc'):
                t1 = time.time()
                x = self._doc_parse(filePath)
                t2 = time.time()

                #print("Time for .doc file " + filePath + ": " + str(t2 - t1))
                
                return x
            elif (fileType == 'docx'):
                t1 = time.time()
                x = self._docx_parse(filePath)
                t2 = time.time()

                #print("Time for .docx file " + filePath + ": " + str(t2 - t1))
                
                return x
            elif (fileType == 'pptx'):
                t1 = time.time()
                x = self._pptx_parse(filePath)
                t2 = time.time()

                #print("Time for .pptx file " + filePath + ": " + str(t2 - t1))
                
                return x
            elif (fileType == 'xlsx'):
                t1 = time.time()
                x = self._xlsx_parse(filePath)
                t2 = time.time()

                #print("Time for .Xlsx file " + filePath + ": " + str(t2 - t1))
                
                return x
            elif (fileType == 'pdf'):
                t1 = time.time()
                x = self._pdf_parse(filePath)
                t2 = time.time()

                #print("Time for .pdf file " + filePath + ": " + str(t2 - t1))
                
                return x
            else:
                t1 = time.time()
                x = self._nonsupported_parse(filePath)
                t2 = time.time()

                #print("Time for other file " + filePath + ": " + str(t2 - t1))
                
                return x

        list = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            futureDownloads = {executor.submit(download_and_parse, file[0], file[1], file[2]): file for file in fileParseType}
            for future in concurrent.futures.as_completed(futureDownloads):
                meta = futureDownloads[future]
                try:
                    return_data = future.result()
                    data = (return_data, meta[2])
                    list.append(data)
                except Exception as exc:
                    print('%r generated an exception: %s' % (meta, exc))
        
        total2 = time.time()

        print("Total parse time for " + str(len(fileList)) + " files: " + str(total2 - total1))

        keywords = ' '.join(search.keywords)

        #Resorts list to follow the original index, for BagOfWords
        list.sort(key=lambda tup: tup[1])

        dataList = []
        for element in list:
            dataList.append(element[0])

        return BagOfWords.find_accurate_docs(fileList, dataList, keywords)

    def __init__(self, dropboxBot): 
        self.dropboxBot = dropboxBot