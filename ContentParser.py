import dropbox
from dropbox.exceptions import AuthError

from io import BytesIO
import tempfile
from docx import Document
from pptx import Presentation
import openpyxl
import PyPDF2
from subprocess import check_output
from SearchFeedback import SearchFeedback
import threading

import time

import concurrent.futures

class ContentParser:
    """
    Class used to parse a list of files' content using multithreads for speed
    -----
    Methods
    -----
    _docx_parse(self, filename)
        Parses through Docx Files
    _doc_parse(self, filename)
        Parses through Doc Files
    _pptx_parse(self, filename)
        Parses through Pptx Files
    _xlsx_parse(self, filename)
        Parses through Xlsx Files
    _pdf_parse(self, filename)
        Parses through Pdf Files
    _nonsupported_parse(self, filename)
        Takes care of nonsupported files
    _determine_parser(fileType, filePath, index)
        Determines which type of parser to use for a file
    parse_file_list(futureParsedList)
        Parses a list of files using inner functionality of the class
    """

    dropboxBot = None

    def _docx_parse(self, filename):
        """
    Parses through Docx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    '\n'.join(fullText)
        Full text of the parsed file
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        doc = Document(stream)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text.lower())

        return '\n'.join(fullText)

    def _doc_parse(self, filename):
        """
    Parses through Doc Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    out.lower()
        the output of the parsed file in lowercase
    """
        out = ''
        with tempfile.NamedTemporaryFile(delete=False, suffix='.doc') as temp:
            self.dropboxBot.dbx.files_download_to_file(temp.name, filename)
            out = check_output(['antiword', '-f' ,temp.name])
        
        return out.lower()

    def _pptx_parse(self, filename):
        """
    Parses through Pptx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    text
        the output of the parsed file
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        #pptx only works with .pptx but not with .ppt
        prs = Presentation(stream)
        textRuns = []
            
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        textRuns.append(run.text.lower())

        text = ''
        for run in textRuns:
            text += run
        
        return text

    def _xlsx_parse(self, filename):
        """
    Parses through Xlsx Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    string
        the output of the parsed file in lowercase
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        stream = BytesIO(f.content)
        book = openpyxl.load_workbook(stream)

        string = ''

        for sheet in book:
            for row in sheet.iter_rows():
                for cell in row:
                    string += ' '
                    string += str(cell.value).lower()

        return string

    def _pdf_parse(self, filename):
        """
    Parses through pdf Files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    string
        the output of the parsed file in lowercase
    """
        try:
            metadata, f = self.dropboxBot.dbx.files_download(filename)
        except dropbox.files.DownloadError as err:
            print(err)
            exit()

        stream = BytesIO(f.content)

        pdfReader = PyPDF2.PdfFileReader(stream)

        numberOfPages = pdfReader.numPages
        string = ''

        for i in range(0, numberOfPages):
            string += ' '
            string += str(pdfReader.getPage(i).extractText().lower())

        return string

    def _nonsupported_parse(self, filename):
        """
    Takes care of nonsupported files

    Parameters
    ----------
    fileName : string
        file name to parse through
    Returns
    -------
    fileName
        the file name
    """
        splitEntry = filename.split('.')
        fileName = ""

        for word in splitEntry[:-1]: # there could be more than one '.'
            fileName += word
        
        return fileName


    def _determine_parser(self, fileType, filePath, index):
        """
        Determines how to parse a given file based on its type that is shown

        Parameters
        ----------
        fileType : string
            the type of file (extension)
        filePath : string
            the path of the file so it can be downloaded for parsing
        index : int
            the index to maintain the position within multiple threads. Not used within this class, 
            but is passed to keep the order for BagOfWords
        Returns
        -------
        content : string
            the content of a file
        """

        if (fileType == 'doc'):
            t1 = time.time()
            content = self._doc_parse(filePath)
            t2 = time.time()
            self.docTime += (t2 - t1)

            #print("Time for .doc file " + filePath + ": " + str(t2 - t1))
            
            return content
        elif (fileType == 'docx'):
            t1 = time.time()
            content = self._docx_parse(filePath)
            t2 = time.time()
            self.docxTime += (t2 - t1)

            #print("Time for .docx file " + filePath + ": " + str(t2 - t1))
            
            return content
        elif (fileType == 'pptx'):
            t1 = time.time()
            content = self._pptx_parse(filePath)
            t2 = time.time()
            self.pptxTime += (t2 - t1)

            #print("Time for .pptx file " + filePath + ": " + str(t2 - t1))
            
            return content
        elif (fileType == 'xlsx'):
            t1 = time.time()
            content = self._xlsx_parse(filePath)
            t2 = time.time()
            self.xlsxTime += (t2 - t1)

            #print("Time for .Xlsx file " + filePath + ": " + str(t2 - t1))
            
            return content
        elif (fileType == 'pdf'):
            t1 = time.time()
            content = self._pdf_parse(filePath)
            t2 = time.time()
            self.pdfTime += (t2 - t1)

            #print("Time for .pdf file " + filePath + ": " + str(t2 - t1))
            
            return content
        else:
            t1 = time.time()
            content = self._nonsupported_parse(filePath)
            t2 = time.time()
            self.nonsupportedTime += (t2 - t1)

            #print("Time for other file " + filePath + ": " + str(t2 - t1))
            
            return content


    def parse_file_list(self, futureParsedList, slackBot, m):
        """
        Parses a list of files using multithreading for speed

        Parameters
        ----------
        futureParsedList : tuple
            a list of tuples containing the type of file (extension), the file path for downloading the file,
            and the index to maintaing the position of the file within concurrency
        Returns
        -------
        list : string
            the list of each file's content to be used for finding the accuracy of the files to the query of the user
        """

        #SearchFeedback.search_feedback(self, len(futureParsedList))
        
        self.numberParsed = 0
        threading.Thread(target=SearchFeedback.search_feedback, args=(self, len(futureParsedList), slackBot, m)).start()

        list = []
        
        self.docTime = 0
        self.docxTime = 0
        self.pptxTime = 0
        self.xlsxTime = 0
        self.pdfTime = 0
        self.nonsupportedTime = 0
        totalTime1 = time.time()

        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
            futureDownloads = {executor.submit(self._determine_parser, file[0], file[1], file[2]): file for file in futureParsedList}
            for future in concurrent.futures.as_completed(futureDownloads):
                meta = futureDownloads[future]
                self.numberParsed += 1
                try:
                    return_data = future.result()
                    data = (return_data, meta[2])
                    list.append(data)
                except Exception as exc:
                    # this is so FileSearches will maintain the same number of files when comparing name searches and content searches
                    emptyData = ("", meta[2]) 
                    list.append(emptyData)
                    #print('%r generated an exception: %s' % (meta, exc))

        print()
        print()
        print("TOTAL TIME TO PARSE DOC FILES: " + str(self.docTime))
        print("TOTAL TIME TO PARSE DOCX FILES: " + str(self.docxTime))
        print("TOTAL TIME TO PARSE PPTX FILES: " + str(self.pptxTime))
        print("TOTAL TIME TO PARSE XLSX FILES: " + str(self.xlsxTime))
        print("TOTAL TIME TO PARSE PDF FILES: " + str(self.pdfTime))
        print("TOTAL TIME TO PARSE NONSUPPORTED FILES: " + str(self.nonsupportedTime))
        print("_______________________________")
        totalTime2 = time.time()
        print("TOTAL TIME TO PARSE " + str(len(futureParsedList)) + " FILES: " + str(totalTime2 - totalTime1))

        return list


    
    def __init__(self, dropboxBot): 
        self.dropboxBot = dropboxBot